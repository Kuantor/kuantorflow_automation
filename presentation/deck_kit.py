#!/usr/bin/env python3
"""KuantorFlow presentation deck kit — reusable python-pptx helpers.

Import `Deck` to build a branded 16:9 slide deck quickly: a consistent
blue/gold theme, a handful of tested layouts (title, content rows, code,
full-bleed image, closing), and speaker-notes support. This is the reusable
core distilled from a full presentation build — the theme and the layout
helpers, without any one deck's hard-coded content.

Run it directly to generate `deck_kit_demo.pptx`, a small sample deck that
exercises every layout — handy as living documentation and a smoke test.

Requires: python-pptx, Pillow (Pillow only for image sizing in `pic_fit`).

    from deck_kit import Deck
    d = Deck()
    d.title_slide("My Talk", "A subtitle", eyebrow="SERIES NAME")
    s = d.content_slide("Section", "Three things")
    d.rows(s, [("First", "why it matters"), ("Second", "..."), ("Third", "...")])
    d.notes(s, "What I'll say over this slide.")
    d.save("out.pptx")
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

try:
    from PIL import Image
except ImportError:  # only needed for pic_fit's aspect-ratio math
    Image = None

# ---- palette (KuantorFlow blue / gold) ----
NAVY   = RGBColor(0x14, 0x23, 0x2E)
NAVY2  = RGBColor(0x1E, 0x33, 0x42)
BLUE   = RGBColor(0x2E, 0x6B, 0xA0)
BLUEL  = RGBColor(0x6F, 0x9F, 0xC8)
GOLD   = RGBColor(0xD9, 0xB1, 0x3C)
INK    = RGBColor(0x22, 0x30, 0x3C)
MUTED  = RGBColor(0x5B, 0x6B, 0x7A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARD   = RGBColor(0xF2, 0xF6, 0xFA)
CARDLN = RGBColor(0xDD, 0xE6, 0xEF)
CODEBG = RGBColor(0x14, 0x23, 0x2E)
CODEFG = RGBColor(0xE8, 0xEF, 0xF5)

HEAD = "Calibri"
BODY = "Calibri"
MONO = "Consolas"
ML = 0.9  # default left margin (inches)


class Deck:
    """A branded 16:9 presentation you build layout-by-layout."""

    def __init__(self, width_in=13.333, height_in=7.5):
        self.prs = Presentation()
        self.prs.slide_width = Inches(width_in)
        self.prs.slide_height = Inches(height_in)
        self.blank = self.prs.slide_layouts[6]
        self.W = self.prs.slide_width
        self.H = self.prs.slide_height

    # ---------- low-level ----------
    def slide(self, bg=WHITE):
        s = self.prs.slides.add_slide(self.blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg
        return s

    def box(self, s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
        tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        return tf

    def para(self, tf, text, size, color, bold=False, font=BODY,
             align=PP_ALIGN.LEFT, before=0, after=4, italic=False,
             first=False, line=None):
        p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(before)
        p.space_after = Pt(after)
        if line:
            p.line_spacing = line
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.name = font
        f.color.rgb = color
        return p

    def rect(self, s, l, t, w, h, fill, line=None, line_w=1,
             shape=MSO_SHAPE.RECTANGLE):
        sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
        sp.shadow.inherit = False
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid()
            sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(line_w)
        return sp

    def pic_fit(self, s, path, l, t, w, h):
        """Place an image contained (and centered) inside a box, in inches."""
        if Image is None:
            raise RuntimeError("Pillow is required for pic_fit(); pip install Pillow")
        iw, ih = Image.open(path).size
        ar = iw / ih
        if ar > w / h:
            pw, ph = w, w / ar
        else:
            ph, pw = h, h * ar
        s.shapes.add_picture(path, Inches(l + (w - pw) / 2), Inches(t + (h - ph) / 2),
                             Inches(pw), Inches(ph))

    def full_bleed(self, s, path):
        s.shapes.add_picture(path, 0, 0, width=self.W, height=self.H)

    def notes(self, s, text):
        s.notes_slide.notes_text_frame.text = text or ""

    # ---------- high-level layouts ----------
    def title_slide(self, title, subtitle=None, eyebrow=None, foot=None, bg=NAVY):
        s = self.slide(bg)
        self.rect(s, 0, 0, self.W.inches, 0.28, GOLD)
        if eyebrow:
            self.para(self.box(s, ML, 2.35, 11.5, 0.5), eyebrow.upper(), 15, GOLD,
                      bold=True, first=True)
        self.para(self.box(s, ML, 2.95, 11.5, 2.0), title, 50, WHITE, bold=True,
                  font=HEAD, first=True, line=1.0)
        if subtitle:
            self.para(self.box(s, ML, 4.75, 11.0, 1.0), subtitle, 22, BLUEL, first=True)
        if foot:
            self.para(self.box(s, ML, 6.7, 11.5, 0.5), foot, 13, MUTED, first=True)
        return s

    def content_slide(self, eyebrow, title, title_size=33):
        """A white content slide with an eyebrow + title drawn; returns it."""
        s = self.slide(WHITE)
        self.para(self.box(s, ML, 0.62, 11.5, 0.4), eyebrow.upper(), 13, BLUE,
                  bold=True, first=True)
        self.para(self.box(s, ML, 1.02, 11.5, 1.1), title, title_size, INK,
                  bold=True, font=HEAD, first=True, line=1.02)
        return s

    def rows(self, s, items, l=ML, t=2.25, w=11.5, gap=0.14, hsize=15,
             dsize=12.5, head_color=BLUE, rowh=None):
        """Stacked (header, description) rows — the workhorse content layout."""
        rh = rowh or min(1.15, 4.7 / max(1, len(items)))
        y = t
        for head, desc in items:
            tf = self.box(s, l, y, w, rh)
            self.para(tf, head, hsize, head_color, bold=True, first=True, after=2, line=1.0)
            if desc:
                self.para(tf, desc, dsize, MUTED, after=0, line=1.02)
            y += rh + gap

    def code_slide(self, eyebrow, title, code, point=None):
        """A dark monospace code panel with an optional highlighted takeaway."""
        s = self.content_slide(eyebrow, title, title_size=30)
        self.rect(s, ML, 2.0, 11.5, 3.7, CODEBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = self.box(s, ML + 0.35, 2.25, 10.9, 3.25)
        for i, ln in enumerate(code.split("\n")):
            self.para(tf, ln if ln else " ", 12.5, CODEFG, font=MONO,
                      first=(i == 0), after=0, line=1.02)
        if point:
            self.rect(s, ML, 6.05, 11.5, 0.95, CARD, line=CARDLN,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            self.para(self.box(s, ML + 0.35, 6.05, 10.9, 0.95, anchor=MSO_ANCHOR.MIDDLE),
                      point, 15, BLUE, bold=True, first=True, line=1.05)
        return s

    def closing_slide(self, title, subtitle=None, bullets=None, footer=None, eyebrow=None):
        s = self.slide(NAVY)
        self.rect(s, 0, 7.22, self.W.inches, 0.28, GOLD)
        if eyebrow:
            self.para(self.box(s, ML, 0.9, 11.5, 0.4), eyebrow.upper(), 13, GOLD,
                      bold=True, first=True)
        self.para(self.box(s, ML, 1.45, 11.5, 1.6), title, 44, WHITE, bold=True,
                  font=HEAD, first=True, line=1.0)
        if subtitle:
            self.para(self.box(s, ML, 2.9, 11.6, 1.5), subtitle, 18, BLUEL,
                      first=True, line=1.18)
        for i, tx in enumerate(bullets or []):
            tf = self.box(s, ML, 4.35 + i * 0.62, 11.5, 0.6)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = "▸  "
            r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = BODY
            r = p.add_run(); r.text = tx
            r.font.size = Pt(15); r.font.color.rgb = RGBColor(0xD3, 0xE0, 0xEC); r.font.name = BODY
        if footer:
            self.para(self.box(s, ML, 6.45, 11.5, 0.7), footer, 16, WHITE, bold=True, first=True)
        return s

    def save(self, path):
        self.prs.save(path)
        return path


def _demo():
    d = Deck()
    d.title_slide("Deck Kit", "A reusable python-pptx starter",
                  eyebrow="KuantorFlow presentations",
                  foot="Run deck_kit.py to regenerate this sample")
    s = d.content_slide("Layouts", "The content-rows layout")
    d.rows(s, [
        ("Title slide", "Dark, branded, gold rule."),
        ("Content rows", "Header + description pairs — the workhorse."),
        ("Code slide", "Dark monospace panel + a takeaway line."),
        ("Closing slide", "Dark recap with gold bullets."),
    ])
    d.notes(s, "Every layout here is a method on Deck.")
    d.code_slide("Example", "A code slide",
                 "def hello(name):\n    return f\"Hello, {name}!\"",
                 "Monospace code with a one-line takeaway.")
    d.closing_slide("Thanks!", subtitle="Import Deck and build your own.",
                    bullets=["One theme, several layouts.",
                             "Speaker notes via d.notes(slide, text)."],
                    footer="deck_kit.py", eyebrow="In one slide")
    out = d.save("deck_kit_demo.pptx")
    print(f"written: {out}  ({len(d.prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    _demo()
